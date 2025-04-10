from fastapi import APIRouter, UploadFile, HTTPException, status
from api.v1.schemas.response_models import (
    StandardResponse, ErrorResponse, ErrorData, SuccessResponse)
from docx import Document
from pptx import Presentation
from starlette.responses import StreamingResponse

import os

word_to_ppt = APIRouter(
    prefix="/word_to_ppt",
    tags=["Word to PPT"],
)
@word_to_ppt.post(
    "/convert",
    response_model=StandardResponse,
    responses={
        200: {
            "description": "Conversion successful",
            "model": StandardResponse,
        },
        400: {
            "description": "Bad Request",
            "model": StandardResponse,
        },
        500: {
            "description": "Internal Server Error",
            "model": StandardResponse,
        },
    },
)
async def convert_to_ppt(file: UploadFile):
    """
    Convert a Word document to a PowerPoint presentation.

    Args:
        file (UploadFile): The Word document to convert.

    Returns:
        StandardResponse: The response containing the converted PPT file.
    """
    try:
        if not file:
            return ErrorResponse(
                status_code=400,
                message="No file provided.",
                data=ErrorData(
                    error="File not found",
                    error_type="Bad Request"
                )
            )
        if not file.filename.endswith('.docx'):
            return ErrorResponse(
                status_code=400,
                message="Invalid file type. Only .docx files are supported.",
                data=ErrorData(
                    error="Invalid file type",
                    error_type="Bad Request"
                )
            )
        doc = Document(file.file)

        ppt = Presentation()

        # Map each heading/paragraph to a slide
        for para in doc.paragraphs:
            slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(slide_layout)
            slide.shapes.title.text = para.text
        else:
            slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Details"
            slide.placeholders[1].text = para.text
        
        # Make ppt directory in root directory
        BASE_DIR = os.path.dirname(os.path.abspath(__file__))
        stripped_path = BASE_DIR.partition('api')[0]
        ppt_dir = os.path.join(stripped_path, "ppt")

        os.makedirs(ppt_dir, exist_ok=True)
        filename = file.filename.strip(".docx")
        # Save the ppt file
        ppt.save(os.path.join(ppt_dir, f"{filename}.pptx"))

        return SuccessResponse(
            status_code=200,
            message="Conversion successful.",
            data={
                "ppt_file": os.path.join(ppt_dir, f"{filename}.pptx")
            }
        )
    except Exception as e:
        return ErrorResponse(
            status_code=500,
            message="An error occurred during conversion.",
            data=ErrorData(
                error=str(e),
                error_type="Internal Server Error"
            )
        )


@word_to_ppt.get(
    "/download/{file_name}",
    )
async def download_file(file_name: str):
    try:
        BASE_DIR = os.path.dirname(os.path.abspath(__file__))
        stripped_path = BASE_DIR.partition('api')[0]
        ppt_dir = os.path.join(stripped_path, "ppt")

        file_path = os.path.join(ppt_dir, f"{file_name}.pptx")

        if not os.path.exists(file_path):
          raise HTTPException(
              status_code=status.HTTP_404_NOT_FOUND,
              detail="File not found"
          )
        
        file_stream = open(file_path, mode="rb")

        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename={file_name}.pptx"
            }
        )
    except Exception as e:
        return ErrorResponse(
            status_code=500,
            message="An error occurred during conversion.",
            data=ErrorData(
                error=str(e),
                error_type="Internal Server Error"
            )
        )